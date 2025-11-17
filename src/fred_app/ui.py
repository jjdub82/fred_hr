# src/fred_app/ui.py
import re
import textwrap
import threading
import tempfile
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import tkinter.font as tkfont

import requests
import numpy as np
import pandas as pd
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk

from .api import search_series, get_series_observations
from .charts import observations_to_df  # expects DataFrame with ["date","value"]
from .config import get_fred_api_key

# -------------------- Preset Report Combos --------------------
PRESET_COMBOS = {
    "Labor Market Conditions": {
        "series": ["UNRATE", "JTSJOL", "CIVPART"],
        "yoy": False, "index": False,
        "dual": {"enabled": True, "series": "JTSJOL"},
    },
    "Wage vs Prices": {
        "series": ["CES0500000003", "ECIWAG", "CPIAUCSL"],
        "yoy": True, "index": False,
        "dual": {"enabled": True, "series": "CPIAUCSL"},
    },
    "Productivity & Compensation": {
        "series": ["OPHNFB", "COMPRNFB", "ULCNFB"],
        "yoy": True, "index": False,
        "dual": {"enabled": True, "series": "ULCNFB"},
    },
    "Hours & Overtime (Mfg)": {
        "series": ["AWHMAN", "AWOTMAN", "TCU"],
        "yoy": False, "index": False,
        "dual": {"enabled": True, "series": "TCU"},
    },
    "Private Payrolls & Temp": {
        "series": ["USPRIV", "TEMPHEL", "JTSHIR"],
        "yoy": True, "index": False,
        "dual": {"enabled": False, "series": "(none)"},
    },
    "Claims & Slack": {
        "series": ["ICSA", "CCSA", "UNRATE"],
        "yoy": False, "index": False,
        "dual": {"enabled": True, "series": "ICSA"},
    },
    "Cost of Living": {
        "series": ["CUSR0000SAH1", "MORTGAGE30US", "USSTHPI"],
        "yoy": True, "index": False,
        "dual": {"enabled": True, "series": "MORTGAGE30US"},
    },
    "Turnover & Retention (JOLTS)": {
        "series": ["JTSQUR", "JTSHIR", "JTSLDR"],
        "yoy": False, "index": False,
        "dual": {"enabled": False, "series": "(none)"},
    },
}

DATE_RE = re.compile(r"^\d{4}-\d{2}-\d{2}$")
FRED_BASE = "https://api.stlouisfed.org/fred"

# -------------------- Styling --------------------
def _apply_ttk_style(root: tk.Tk):
    """Readable light theme; no dark surprises."""
    try:
        import sv_ttk
        sv_ttk.set_theme("light")
    except Exception:
        style = ttk.Style(root)
        try: style.theme_use("clam")
        except Exception: pass
        style.configure(".", font=("Segoe UI", 10))
        style.configure("TButton", padding=6)
        style.configure("TEntry", padding=4)
        style.configure("TCombobox", padding=4)
        style.configure("Treeview.Heading", padding=6)
        style.map("Treeview",
                  background=[("selected", "#2a6fdb")],
                  foreground=[("selected", "white")])

def _apply_mpl_style():
    try: plt.style.use("seaborn-v0_8-whitegrid")
    except Exception: pass
    mpl.rcParams.update({
        "figure.autolayout": True,
        "axes.spines.top": False, "axes.spines.right": False,
        "axes.grid": True, "grid.alpha": 0.25,
        "lines.linewidth": 2.0, "font.size": 10
    })

# -------------------- FRED helpers --------------------
def _fred_get(path: str, **params):
    p = {**params, "api_key": get_fred_api_key(), "file_type": "json"}
    r = requests.get(f"{FRED_BASE}/{path}", params=p, timeout=20)
    r.raise_for_status()
    return r.json()

def fred_top_categories():
    return _fred_get("category/children", category_id=0).get("categories", [])

def fred_children(category_id: int):
    return _fred_get("category/children", category_id=category_id).get("categories", [])

def fred_category_series(category_id: int, limit: int = 1000):
    return _fred_get("category/series",
                     category_id=category_id,
                     order_by="popularity",
                     sort_order="desc",
                     limit=limit).get("seriess", [])

def fred_series_title(series_id: str) -> str:
    """Fetch the official FRED title for a series id."""
    try:
        js = _fred_get("series", series_id=series_id)
        items = js.get("seriess", [])
        if items:
            return str(items[0].get("title", "")).strip() or series_id
    except Exception:
        pass
    return series_id


# -------------------- Main Window --------------------
class MainWindow(ttk.Frame):
    """FRED desktop UI (light theme)."""

    def __init__(self, master: tk.Tk):
        _apply_ttk_style(master)
        _apply_mpl_style()
        super().__init__(master, padding=12)
        self.master.title("FRED App — HR Edition (Light)")
        self.master.geometry("1400x780")
        self.master.minsize(1120, 680)

        # state
        self.selected_ids: list[str] = []
        self.series_labels: dict[str, str] = {}
        self.current_df: pd.DataFrame | None = None
        self.canvas: FigureCanvasTkAgg | None = None
        self.toolbar: NavigationToolbar2Tk | None = None

        # category state
        self.cat_map: dict[str, int] = {}
        self.subcat_map: dict[str, int] = {}

        # scenario / dual-axis state
        self.scenario_var = tk.StringVar(master=self, value="Baseline")
        self.dual_axis_var = tk.BooleanVar(master=self, value=False)
        self.dual_series_var = tk.StringVar(master=self, value="(none)")

        self._build_layout()
        self.pack(fill="both", expand=True)

        # bindings / startup
        self.master.bind("<Control-Return>", lambda e: self._run_query())
        self.master.bind("<F5>", lambda e: self._run_query())
        self._load_categories_bg()
        self.master.protocol("WM_DELETE_WINDOW", self._on_close)

    # -------------------- Layout --------------------
    def _build_layout(self):
        # -------- Top row: search / range / toggles --------
        top1 = ttk.Frame(self); top1.pack(fill="x", pady=(0, 6))
        ttk.Label(top1, text="Search").pack(side="left", padx=(0, 6))
        self.q = tk.StringVar()
        ent_q = ttk.Entry(top1, textvariable=self.q, width=36); ent_q.pack(side="left")
        ent_q.bind("<Return>", lambda e: self._on_search())
        ttk.Button(top1, text="Go", command=self._on_search).pack(side="left", padx=6)

        ttk.Label(top1, text="Range").pack(side="left", padx=(16, 6))
        self.start_var = tk.StringVar(value=""); self.end_var = tk.StringVar(value="")
        ent_s = ttk.Entry(top1, textvariable=self.start_var, width=12); ent_s.pack(side="left")
        ttk.Label(top1, text="to").pack(side="left", padx=4)
        ent_e = ttk.Entry(top1, textvariable=self.end_var, width=12); ent_e.pack(side="left")
        ent_s.bind("<Return>", lambda e: self._run_query()); ent_e.bind("<Return>", lambda e: self._run_query())

        ttk.Label(top1, text="Quick Range").pack(side="left", padx=(16, 6))
        self.range_var = tk.StringVar(value="Max")
        quick = ttk.Combobox(top1, textvariable=self.range_var, values=["YTD","1Y","5Y","10Y","Max"],
                             width=6, state="readonly")
        quick.pack(side="left"); quick.bind("<<ComboboxSelected>>", lambda e: self._apply_quick_range())

        self.normalize_var = tk.BooleanVar(master=self, value=False)
        self.yoy_var = tk.BooleanVar(master=self, value=False)
        ttk.Checkbutton(top1, text="Index=100", variable=self.normalize_var, command=self._refresh_chart)\
            .pack(side="left", padx=(16, 4))
        ttk.Checkbutton(top1, text="YoY %", variable=self.yoy_var, command=self._refresh_chart).pack(side="left")

        self.forecast_var = tk.BooleanVar(master=self, value=False)
        ttk.Checkbutton(top1, text="Trend range", variable=self.forecast_var, command=self._refresh_chart)\
            .pack(side="left", padx=(16, 0))

        # -------- Second row: Combos (first), Dual, Horizon/Scenario, Actions --------
        top2 = ttk.Frame(self); top2.pack(fill="x", pady=(0, 8))

        left_controls = ttk.Frame(top2); left_controls.pack(side="left")

        # --- Combos block (always visible; auto-add on selection) ---
        combo_wrap = ttk.Frame(left_controls); combo_wrap.pack(side="left", padx=(0, 12))
        ttk.Label(combo_wrap, text="Combos").pack(side="left", padx=(0, 6))
        self.combo_var = tk.StringVar(value="(choose)")
        self.combo_cb = ttk.Combobox(combo_wrap,
                                     textvariable=self.combo_var,
                                     values=sorted(PRESET_COMBOS.keys()),
                                     state="readonly")
        self._autosize_combo_width()  # prevent truncation
        self.combo_cb.pack(side="left")
        self.combo_cb.bind("<<ComboboxSelected>>", lambda e: self._add_combo(auto=True))
        ttk.Button(combo_wrap, text="Add Combo", command=self._add_combo).pack(side="left", padx=(6, 0))

        # Dual axis + selector
        ttk.Checkbutton(left_controls, text="Dual axis",
                        variable=self.dual_axis_var, command=self._on_dual_toggle).pack(side="left", padx=(8, 6))
        ttk.Label(left_controls, text="2nd metric").pack(side="left", padx=(0, 4))
        self.dual_series_cb = ttk.Combobox(left_controls, textvariable=self.dual_series_var,
                                           values=["(none)"], width=28, state="disabled")
        self.dual_series_cb.pack(side="left")
        self.dual_series_cb.bind("<<ComboboxSelected>>", lambda e: self._refresh_chart())

        # Horizon + Scenario
        ttk.Label(left_controls, text="Horizon").pack(side="left", padx=(12, 4))
        self.forecast_h_var = tk.StringVar(master=self, value="12m")
        h_cb = ttk.Combobox(left_controls, textvariable=self.forecast_h_var,
                            values=["6m","12m","18m","24m"], width=5, state="readonly")
        h_cb.pack(side="left", padx=(0, 12)); h_cb.bind("<<ComboboxSelected>>", lambda e: self._refresh_chart())

        ttk.Label(left_controls, text="Scenario").pack(side="left", padx=(0, 4))
        scen_cb = ttk.Combobox(left_controls, textvariable=self.scenario_var,
                               values=["Baseline","Hot labor","Cooldown"], width=12, state="readonly")
        scen_cb.pack(side="left", padx=(0, 12)); scen_cb.bind("<<ComboboxSelected>>", lambda e: self._refresh_chart())

        # Right actions (trimmed for space)
        actions = ttk.Frame(top2); actions.pack(side="right")
        ttk.Button(actions, text="Export CSV", command=self._export_csv).pack(side="right", padx=(6, 0))
        ttk.Button(actions, text="Export PPT", command=self._export_ppt).pack(side="right", padx=(6, 0))

        # -------- Category/Subcategory bar --------
        catbar = ttk.Frame(self); catbar.pack(fill="x", pady=(0, 8))
        ttk.Label(catbar, text="Category").pack(side="left")
        self.category_var = tk.StringVar(value="")
        self.category_cb = ttk.Combobox(catbar, textvariable=self.category_var, width=40, state="readonly")
        self.category_cb.pack(side="left", padx=(6, 12))
        self.category_cb.bind("<<ComboboxSelected>>", lambda e: self._on_category_selected())

        ttk.Label(catbar, text="Subcategory").pack(side="left")
        self.subcategory_var = tk.StringVar(value="")
        self.subcategory_cb = ttk.Combobox(catbar, textvariable=self.subcategory_var, width=40, state="readonly")
        self.subcategory_cb.pack(side="left", padx=(6, 12))
        self.subcategory_cb.bind("<<ComboboxSelected>>", lambda e: self._on_subcategory_selected())

        ttk.Button(catbar, text="Load Series", command=self._load_series_from_selection).pack(side="left")
        ttk.Button(catbar, text="Clear Filters", command=self._clear_filters).pack(side="left", padx=(8, 0))

        # -------- Body split (LEFT | RIGHT) --------
        body = ttk.Panedwindow(self, orient=tk.HORIZONTAL); body.pack(fill="both", expand=True)

        # ------ LEFT (vertical split; RUN always visible) ------
        left = ttk.Frame(body, padding=(0, 0, 10, 0)); body.add(left, weight=1)

        left_split = tk.PanedWindow(left, orient=tk.VERTICAL, sashrelief="flat", bd=0)
        left_split.pack(fill="both", expand=True)

        # --- Top pane: Results list ---
        res_wrap = ttk.Frame(left_split); res_wrap.grid_rowconfigure(0, weight=1); res_wrap.grid_columnconfigure(0, weight=1)
        res_frame = ttk.Frame(res_wrap); res_frame.grid(row=0, column=0, sticky="nsew", pady=(0, 6))
        res_frame.grid_rowconfigure(0, weight=1); res_frame.grid_columnconfigure(0, weight=1)

        self.results_tree = ttk.Treeview(res_frame, columns=("id", "title"), show="headings", style="Results.Treeview")
        self.results_tree.heading("id", text="Series ID"); self.results_tree.heading("title", text="Title")
        self.results_tree.column("id", width=180, anchor="w"); self.results_tree.column("title", width=460, anchor="w")
        res_vsb = ttk.Scrollbar(res_frame, orient="vertical", command=self.results_tree.yview)
        self.results_tree.configure(yscrollcommand=res_vsb.set)
        self.results_tree.grid(row=0, column=0, sticky="nsew"); res_vsb.grid(row=0, column=1, sticky="ns")

        left_split.add(res_wrap, minsize=160)

        # --- Bottom pane: Selected Series + RUN ---
        sel_wrap = ttk.Frame(left_split)

        btns = ttk.Frame(sel_wrap); btns.pack(fill="x", pady=(0, 6))
        left_btns = ttk.Frame(btns); left_btns.pack(side="left")
        ttk.Button(left_btns, text="Add →", command=self._add_from_results).pack(side="left")
        ttk.Button(left_btns, text="← Remove", command=self._remove_from_selected).pack(side="left", padx=6)
        ttk.Button(left_btns, text="Clear Selected", command=self._clear_selected).pack(side="left", padx=6)
        ttk.Button(btns, text="RUN ▶", command=self._run_query).pack(side="right")

        ttk.Label(sel_wrap, text="Selected Series").pack(anchor="w", pady=(4, 2))
        list_frame = ttk.Frame(sel_wrap); list_frame.pack(fill="both", expand=True)
        list_frame.grid_rowconfigure(0, weight=1); list_frame.grid_columnconfigure(0, weight=1)

        self.sel_list = tk.Listbox(list_frame, height=8)
        sel_vsb = ttk.Scrollbar(list_frame, orient="vertical", command=self.sel_list.yview)
        self.sel_list.configure(yscrollcommand=sel_vsb.set)
        self.sel_list.grid(row=0, column=0, sticky="nsew"); sel_vsb.grid(row=0, column=1, sticky="ns")
        self.sel_list.bind("<Double-1>", lambda e: self._run_query()); self.sel_list.bind("<Return>", lambda e: self._run_query())

        left_split.add(sel_wrap, minsize=200)

        # ------ RIGHT ------
        right = ttk.Notebook(body); body.add(right, weight=3)

        self.chart_tab = ttk.Frame(right, padding=6); right.add(self.chart_tab, text="Chart")
        self.table_tab = ttk.Frame(right, padding=6); right.add(self.table_tab, text="Table")
        self.insights_tab = ttk.Frame(right, padding=10); right.add(self.insights_tab, text="Insights")

        self.insights_text = tk.Text(self.insights_tab, height=12, wrap="word")
        self.insights_text.configure(state="disabled"); self.insights_text.pack(fill="both", expand=True)

        self.chart_frame = ttk.Frame(self.chart_tab); self.chart_frame.pack(fill="both", expand=True)
        self.toolbar_frame = ttk.Frame(self.chart_tab); self.toolbar_frame.pack(fill="x")

        self.table = ttk.Treeview(self.table_tab, show="headings", style="Data.Treeview"); self.table.pack(fill="both", expand=True)

        # Styles + bindings
        self._apply_treeview_colors()
        self.results_tree.bind("<Double-1>", lambda e: self._add_from_results())
        self.results_tree.bind("<Return>", lambda e: self._add_from_results())

        # Status bar
        status = ttk.Frame(self); status.pack(fill="x", pady=(10, 0))
        self.status_var = tk.StringVar(value="Ready"); ttk.Label(status, textvariable=self.status_var, anchor="w").pack(fill="x")

    # --- Treeview styles ---
    def _apply_treeview_colors(self):
        style = ttk.Style(self)
        base_bg = "#FFFFFF"; alt_bg = "#F2F4F7"; base_fg = "#111827"
        sel_bg = "#2A6FDB"; sel_fg = "#FFFFFF"
        style.configure("Results.Treeview", background=base_bg, fieldbackground=base_bg, foreground=base_fg, rowheight=48)
        style.configure("Results.Treeview.Heading", foreground=base_fg)
        style.map("Results.Treeview", background=[("selected", sel_bg)], foreground=[("selected", sel_fg)])
        style.configure("Data.Treeview", background=base_bg, fieldbackground=base_bg, foreground=base_fg, rowheight=28)
        style.configure("Data.Treeview.Heading", foreground=base_fg)
        style.map("Data.Treeview", background=[("selected", sel_bg)], foreground=[("selected", sel_fg)])
        for tv in (getattr(self, "results_tree", None), getattr(self, "table", None)):
            if tv:
                tv.tag_configure("even", background=base_bg, foreground=base_fg)
                tv.tag_configure("odd", background=alt_bg, foreground=base_fg)

    # -------------------- Combos handler --------------------
    def _autosize_combo_width(self):
        """Make the Combos combobox wide enough for the longest preset label."""
        try:
            longest = max((len(k) for k in PRESET_COMBOS.keys()), default=20)
            width_chars = max(28, min(60, longest + 4))
            self.combo_cb.configure(width=width_chars)
        except Exception:
            self.combo_cb.configure(width=40)

    def _add_combo(self, auto: bool = False):
        name = (self.combo_var.get() or "").strip()
        if not name or name not in PRESET_COMBOS:
            if not auto:
                messagebox.showinfo("Info", "Pick a combo from the list.")
            return

        cfg = PRESET_COMBOS[name]
        series = cfg.get("series", [])
        added = 0
        for sid in series:
            if sid not in self.selected_ids:
                self.selected_ids.append(sid)
                self.sel_list.insert("end", sid)
                # NEW: keep a nice label for charts/insights
                self.series_labels[sid] = fred_series_title(sid)
                added += 1


        # Suggested view defaults from combo
        self.yoy_var.set(bool(cfg.get("yoy", False)))
        self.normalize_var.set(bool(cfg.get("index", False)))
        dual_cfg = cfg.get("dual", {}) or {}
        use_dual = bool(dual_cfg.get("enabled", False))
        self.dual_axis_var.set(use_dual)
        pick = dual_cfg.get("series", "(none)") or "(none)"
        self.dual_series_var.set(pick)
        self.dual_series_cb.configure(state=("readonly" if use_dual else "disabled"))

        self._set_status(f"Added {added} series from '{name}'. Press RUN ▶ to fetch data.")

    # -------------------- Categories & Search --------------------
    def _load_categories_bg(self):
        self._set_status("Loading categories…")
        def work(): return fred_top_categories()
        def done(top):
            top_sorted = sorted(top, key=lambda c: c.get("name", "").lower())
            self.cat_map = {c["name"]: c["id"] for c in top_sorted}
            self.category_cb["values"] = list(self.cat_map.keys())
            self._set_status(f"Ready – {len(top_sorted)} categories loaded" if top_sorted else "Ready")
        threading.Thread(target=lambda: self._bg(work, done), daemon=True).start()

    def _on_category_selected(self):
        name = self.category_var.get().strip()
        if not name: return
        cat_id = self.cat_map.get(name); 
        if not cat_id: return
        self._set_status(f"Loading subcategories for “{name}”…")
        def work(): return sorted(fred_children(cat_id), key=lambda c: c.get("name","").lower())
        def done(children_sorted):
            self.subcat_map = {c["name"]: c["id"] for c in children_sorted}
            self.subcategory_cb["values"] = list(self.subcat_map.keys())
            self.subcategory_var.set(""); self._set_status("Ready")
        threading.Thread(target=lambda: self._bg(work, done), daemon=True).start()

    def _on_subcategory_selected(self):
        self._load_series_from_selection()

    def _load_series_from_selection(self):
        subname = self.subcategory_var.get().strip()
        catname = self.category_var.get().strip()
        if subname and subname in self.subcat_map:
            cid, scope = self.subcat_map[subname], subname
        elif catname and catname in self.cat_map:
            cid, scope = self.cat_map[catname], catname
        else:
            messagebox.showinfo("Info", "Choose a Category (and optionally a Subcategory) first.")
            return
        self._set_status(f"Loading series for “{scope}”…")
        def work(): return fred_category_series(cid, limit=1000)
        def done(series_list):
            self._clear_results()
            for i, s in enumerate(series_list):
                tag = "even" if i % 2 == 0 else "odd"
                sid = s.get("id",""); title = self._wrap_title(s.get("title",""), width=50)
                self.results_tree.insert("", "end", values=(sid, title), tags=(tag,))
            self._refresh_results_rowheight(); self._set_status(f"Loaded {len(series_list)} series from “{scope}”")
        threading.Thread(target=lambda: self._bg(work, done), daemon=True).start()

    def _on_search(self):
        term = self.q.get().strip()
        if not term:
            messagebox.showinfo("Info", "Enter a search term."); return
        self._set_status("Searching…")
        def work(): return search_series(term, limit=200)
        def done(results):
            self._clear_results()
            for i, s in enumerate(results):
                tag = "even" if i % 2 == 0 else "odd"
                sid = s.get("id",""); title = self._wrap_title(s.get("title",""), width=50)
                self.results_tree.insert("", "end", values=(sid, title), tags=(tag,))
            self._refresh_results_rowheight(); self._set_status(f"Found {len(results)} result(s)")
        threading.Thread(target=lambda: self._bg(work, done), daemon=True).start()

    def _clear_results(self):
        for r in self.results_tree.get_children(): self.results_tree.delete(r)

    # -------------------- Selection handlers --------------------
    def _add_from_results(self):
        sel = self.results_tree.selection()
        if not sel: self._set_status("Pick a series first."); return
        item = self.results_tree.item(sel[0]); vals = item.get("values") or []
        sid = vals[0] if len(vals)>0 else ""; title = vals[1] if len(vals)>1 else ""
        if not sid: self._set_status("Couldn’t read a series ID."); return
        if sid not in self.selected_ids:
            self.selected_ids.append(sid); self.sel_list.insert("end", sid)
        if title: self.series_labels[sid] = str(title).replace("\n"," ").strip()
        self._set_status(f"Added {sid}")

    def _remove_from_selected(self):
        sel = list(self.sel_list.curselection()); sel.reverse()
        for idx in sel:
            sid = self.sel_list.get(idx); self.sel_list.delete(idx)
            if sid in self.selected_ids: self.selected_ids.remove(sid)

    def _clear_selected(self):
        self.sel_list.delete(0, "end"); self.selected_ids.clear()

    # -------------------- Run & data --------------------
    def _run_query(self):
        if not self.selected_ids:
            messagebox.showinfo("Info", "Select at least one series first.")
            self._set_status("No series selected."); return

        start, end = self._clean_dates(self.start_var.get().strip(), self.end_var.get().strip())
        params = {}
        if start: params["observation_start"] = start
        if end: params["observation_end"] = end

        preview = ", ".join(self.selected_ids[:3]) + ("..." if len(self.selected_ids) > 3 else "")
        self._set_status(f"Fetching series: {preview}")

        def work():
            dfs = []
            for sid in self.selected_ids:
                try:
                    obs = get_series_observations(sid, **params)
                    df = observations_to_df(obs).set_index("date").rename(columns={"value": sid})
                    dfs.append(df)
                except Exception as e:
                    print(f"Error fetching {sid}: {e}")
            return pd.concat(dfs, axis=1) if dfs else pd.DataFrame()

        def done(df: pd.DataFrame):
            self.current_df = df
            opts = ["(none)"] + list(df.columns)
            self.dual_series_cb["values"] = opts
            if self.dual_series_var.get() not in opts: self.dual_series_var.set("(none)")
            self.dual_series_cb.configure(state=("readonly" if self.dual_axis_var.get() else "disabled"))

            if df.empty:
                self._set_status("No data for selection/date range.")
                messagebox.showinfo("No Data", "No observations returned.")
                self._clear_chart(); self._clear_table(); self._refresh_insights(); return

            self._refresh_chart(); self._refresh_table(); self._refresh_insights()
            self._set_status(f"Loaded {len(df.columns)} series, {len(df)} rows.")

        threading.Thread(target=lambda: self._bg(work, done), daemon=True).start()

    # -------------------- Dates / ranges --------------------
    def _clean_dates(self, start: str, end: str):
        valid_s = bool(DATE_RE.match(start)); valid_e = bool(DATE_RE.match(end))
        if not valid_s and start: self._set_status("Ignored invalid start date (YYYY-MM-DD)."); start = ""
        if not valid_e and end: self._set_status("Ignored invalid end date (YYYY-MM-DD)."); end = ""
        if start and end:
            try:
                s, e = pd.to_datetime(start), pd.to_datetime(end)
                if s > e: start, end = end, start; self._set_status("Swapped start/end to maintain order.")
            except Exception:
                start, end = "", ""; self._set_status("Ignored invalid date range.")
        return start, end

    def _apply_quick_range(self):
        choice = self.range_var.get()
        try:
            today = pd.Timestamp.today().normalize()
            if choice == "YTD":
                s = pd.Timestamp(today.year, 1, 1); e = today
            elif choice == "1Y":
                s = today - pd.DateOffset(years=1); e = today
            elif choice == "5Y":
                s = today - pd.DateOffset(years=5); e = today
            elif choice == "10Y":
                s = today - pd.DateOffset(years=10); e = today
            elif choice == "Max":
                self.start_var.set(""); self.end_var.set(""); return
            else: return
            self.start_var.set(s.date().isoformat()); self.end_var.set(e.date().isoformat())
        except Exception: pass

    # -------------------- Forecasting utilities --------------------
    def _infer_freq_and_steps(self, idx: pd.DatetimeIndex):
        if not isinstance(idx, pd.DatetimeIndex) or len(idx) < 3: return None, 0
        f = pd.infer_freq(idx)
        if f in ("M","MS","ME","BM","BMS"): return "M", 24
        if f and str(f).upper().startswith("Q"): return "Q", 8
        return "M", 24

    def _forecast_steps_from_ui(self, freq: str) -> int:
        m = {"6m":6, "12m":12, "18m":18, "24m":24}
        months = m.get(self.forecast_h_var.get(), 12)
        if freq == "Q":
            import math; return max(2, int(math.ceil(months/3)))
        return months

    def _scenario_params(self):
        scen = (self.scenario_var.get() or "Baseline").lower()
        if "hot" in scen: return 0.45, 1.35, -0.10
        if "cool" in scen: return 0.30, 1.10, +0.10
        return 0.35, 1.28, 0.0

    def _trend_projection_band(self, y: pd.Series, steps: int, freq: str, lookback: int = 60):
        y = y.dropna()
        if len(y) < 12: return None
        if len(y) > lookback: y = y.iloc[-lookback:]
        last = y.index[-1]
        idx_future = (pd.date_range(last + pd.offsets.QuarterEnd(1), periods=steps, freq="Q")
                      if freq == "Q" else
                      pd.date_range(last + pd.offsets.MonthEnd(1), periods=steps, freq="M"))
        level_med = float(np.nanmedian(y.values)); level_max = float(np.nanmax(y.values))
        is_rate_like = (level_med < 20.0 and level_max < 50.0)
        diffs = y.diff().dropna(); sigma = float(diffs.std()) if len(diffs) >= 6 else 0.0
        kappa_annual, band_mult_k, mu_shift = self._scenario_params()
        if is_rate_like:
            try:
                mu = float(y.rolling(60, min_periods=12).mean().iloc[-1])
                if np.isnan(mu): mu = float(y.mean())
            except Exception: mu = float(y.mean())
            mu = max(0.0, mu * (1.0 + mu_shift))
            kappa = kappa_annual / (12 if freq == "M" else 4)
            mid = np.empty(steps, dtype=float); x = float(y.iloc[-1])
            for i in range(steps): x = x + kappa * (mu - x); mid[i] = x
        else:
            t = np.arange(len(y), dtype=float)
            a, b = np.polyfit(t, y.values.astype(float), 1)
            mean_level = float(np.nanmean(y.values)); cap = 0.02 * max(mean_level, 1.0)
            b = float(np.clip(b, -cap, cap)); tf = np.arange(len(y), len(y)+steps, dtype=float)
            mid = a + b * tf
        h = np.arange(1, steps+1, dtype=float)
        band = band_mult_k * sigma * np.sqrt(h)
        upper = mid + band; lower = mid - band
        ymax_recent = max(level_max, np.nanmax(y.values))
        cap_hi = max(ymax_recent * 1.5, np.nanmax(mid)); cap_lo = 0.0 if is_rate_like else -np.inf
        mid = np.clip(mid, cap_lo, cap_hi); upper = np.clip(upper, cap_lo, cap_hi); lower = np.clip(lower, cap_lo, cap_hi)
        return pd.DataFrame({"mid": mid, "low": lower, "high": upper}, index=idx_future)

    # -------------------- Chart & Table --------------------
    def _wrap_label(self, s: str, width: int = 32) -> str:
        s = " ".join(str(s or "").split())
        return "\n".join(textwrap.wrap(s, width=width)) if len(s) > width else s

    def _on_dual_toggle(self):
        self.dual_series_cb.configure(state=("readonly" if self.dual_axis_var.get() else "disabled"))
        self._refresh_chart()

    def _refresh_chart(self):
        df = self.current_df
        if df is None or df.empty:
            self._clear_chart(); return
        df = df.copy(); df.index = pd.to_datetime(df.index); df = df.sort_index()
        df = df.apply(pd.to_numeric, errors="coerce")
        plot = df.copy()
        if self.yoy_var.get(): plot = plot.pct_change(12) * 100.0
        if self.normalize_var.get() and not self.yoy_var.get():
            first = plot.dropna(how="all").iloc[0]; plot = plot.divide(first) * 100.0
        plot = plot.astype("float64").dropna(how="all", axis=1)
        if plot.empty:
            self._clear_chart(); self._set_status("No plottable data."); return

        fig = plt.figure(); ax = fig.add_subplot(111)
        labels = []; eligible_cnt = 0

        use_dual = bool(self.dual_axis_var.get())
        dual_target = self.dual_series_var.get()
        use_dual = use_dual and dual_target and dual_target != "(none)" and dual_target in plot.columns
        ax2 = None; dual_line = None
        dual_units_label = "YoY % " if self.yoy_var.get() else "Value"

        for col in plot.columns:
            label = self.series_labels.get(col, col); labels.append(label)
            series = plot[col].astype("float64")
            if use_dual and col == dual_target:
                if ax2 is None: ax2 = ax.twinx()
                dual_line = ax2.plot(plot.index, series.values, linestyle="--", linewidth=2.0,
                                     label=self._wrap_label(f"{label} (2nd axis)", 32))[0]
            else:
                ax.plot(plot.index, series.values, label=self._wrap_label(label, 32))

            if self.forecast_var.get() and not (use_dual and col == dual_target):
                freq, steps_cap = self._infer_freq_and_steps(series.index)
                if freq:
                    steps = min(steps_cap, self._forecast_steps_from_ui(freq))
                    band = self._trend_projection_band(series, steps=steps, freq=freq)
                    if band is not None and not band.empty:
                        x = band.index; mid = band["mid"].to_numpy(dtype=float)
                        low = band["low"].to_numpy(dtype=float); high = band["high"].to_numpy(dtype=float)
                        ax.fill_between(x, low, high, alpha=0.12, linewidth=0)
                        ax.plot(x, mid, linestyle="--", linewidth=2.2,
                                label=self._wrap_label(f"{label} (trend range)", 32))
                        eligible_cnt += 1

        ax.set_title(self._wrap_label(labels[0], 48) if len(labels)==1 else "FRED Series", pad=10)
        ax.set_xlabel("Date"); ax.set_ylabel("YoY % " if self.yoy_var.get() else "Value")
        if use_dual and ax2 is not None:
            ax2.set_ylabel(dual_units_label)
            if dual_line is not None:
                c = dual_line.get_color(); ax2.yaxis.label.set_color(c); ax2.tick_params(axis="y", colors=c)

        if len(labels) <= 2 and not use_dual:
            ax.legend(loc="best", frameon=False, fontsize=9); fig.tight_layout()
        else:
            h1, l1 = ax.get_legend_handles_labels()
            if use_dual and dual_line is not None: h1 += [dual_line]; l1 += [dual_line.get_label()]
            fig.subplots_adjust(right=0.78)
            lg = ax.legend(h1, l1, loc="upper left", bbox_to_anchor=(1.01, 1.0), frameon=False, fontsize=9)
            for t in lg.get_texts(): t.set_ha("left")

        fig.autofmt_xdate(rotation=0, ha="center"); self._set_chart(fig)
        if self.forecast_var.get():
            self._set_status(f"Projected {eligible_cnt} series" if eligible_cnt else "No series eligible for trend range.")

    def _refresh_table(self):
        df = self.current_df
        if df is None or df.empty:
            self._clear_table(); return
        table_df = df.copy()
        last = table_df.index.max(); start = last - pd.DateOffset(months=12)
        snap = table_df.loc[table_df.index >= start].copy()

        self.table.delete(*self.table.get_children())
        self.table["columns"] = ["date"] + list(table_df.columns)
        for c in self.table["columns"]:
            self.table.heading(c, text=c); self.table.column(c, width=120 if c=="date" else 110, anchor="e")

        for i, row in enumerate(snap.reset_index().itertuples(index=False)):
            vals = [row.date.date().isoformat()] + [
                "" if pd.isna(getattr(row, col)) else f"{getattr(row, col):,.3f}" for col in table_df.columns
            ]
            tag = "even" if i % 2 == 0 else "odd"
            self.table.insert("", "end", values=vals, tags=(tag,))

    def _set_chart(self, fig):
        if self.canvas: self.canvas.get_tk_widget().destroy(); self.canvas = None
        if self.toolbar: self.toolbar.destroy(); self.toolbar = None
        self.canvas = FigureCanvasTkAgg(fig, master=self.chart_frame)
        self.canvas.draw(); self.canvas.get_tk_widget().pack(fill="both", expand=True)
        self.toolbar = NavigationToolbar2Tk(self.canvas, self.toolbar_frame, pack_toolbar=False)
        self.toolbar.update(); self.toolbar.pack(side="right", padx=4, pady=2)

    def _clear_chart(self):
        if self.canvas: self.canvas.get_tk_widget().destroy(); self.canvas = None
        if self.toolbar: self.toolbar.destroy(); self.toolbar = None

    def _clear_table(self):
        self.table.delete(*self.table.get_children()); self.table["columns"] = []

    # -------------------- Insights --------------------
    def _refresh_insights(self):
        self.insights_text.configure(state="normal"); self.insights_text.delete("1.0", "end")
        df = self.current_df
        if df is None or df.empty:
            self.insights_text.insert("end", "Run a query to see insights.")
            self.insights_text.configure(state="disabled"); return

        work_df = df.apply(pd.to_numeric, errors="coerce").sort_index()
        last_idx = work_df.dropna(how="all").index.max()
        if pd.isna(last_idx):
            self.insights_text.insert("end", "No recent data to summarize.")
            self.insights_text.configure(state="disabled"); return

        lines = [f"Snapshot date: {pd.to_datetime(last_idx).date().isoformat()}", ""]
        alerts = []
        for col in work_df.columns:
            s = work_df[col].dropna()
            if len(s) < 6: continue
            title = self.series_labels.get(col, col)
            info = self._series_brief_stats(s); tp = self._detect_turning_point(s); zf = self._zscore_flag(s)
            parts = [f"Latest {info['latest']:.2f}"]
            if info["d3"] is not None: parts.append(f"Δ3m {info['d3']:+.2f}")
            if info["d6"] is not None: parts.append(f"Δ6m {info['d6']:+.2f}")
            if info["d12"] is not None: parts.append(f"Δ12m {info['d12']:+.2f}")
            if info["yoy"] is not None: parts.append(f"YoY {info['yoy']:+.2f}%")
            if info["slope_ann"] is not None: parts.append(f"Trend {info['slope_ann']:+.2f}/yr")
            lines.append(f"• {title}\n    " + " | ".join(parts))
            flags = []
            if zf is not None and abs(zf["z"]) >= 2 and zf["window"] >= 18:
                flags.append(f"anomaly (z={zf['z']:+.1f} vs {zf['window']}m mean)")
            if tp is not None and tp["flip"]:
                flags.append(f"turning point ({tp['from']} → {tp['to']})")
            if flags:
                lines.append("    ⚑ " + " & ".join(flags)); alerts.append(f"{title}: " + ", ".join(flags))

        if alerts:
            lines.append(""); lines.append("Alerts:")
            for a in alerts: lines.append(f"    • {a}")

        self.insights_text.insert("end", "\n".join(lines).strip()); self.insights_text.configure(state="disabled")

    def _infer_months_in_year(self, s: pd.Series) -> int:
        if not isinstance(s.index, pd.DatetimeIndex) or len(s) < 3: return 12
        f = pd.infer_freq(s.index); 
        if f and str(f).upper().startswith("Q"): return 4
        return 12

    def _series_brief_stats(self, s: pd.Series) -> dict:
        s = s.dropna()
        out = {"latest": float(s.iloc[-1]), "d3": None, "d6": None, "d12": None, "yoy": None, "slope_ann": None}
        def safe_delta(n):
            if len(s) > n:
                try: return float(s.iloc[-1] - s.iloc[-1-n])
                except Exception: return None
            return None
        out["d3"] = safe_delta(3); out["d6"] = safe_delta(6); out["d12"] = safe_delta(12)
        try:
            yoy = s.pct_change(12).iloc[-1] * 100.0
            if pd.notna(yoy) and np.isfinite(yoy): out["yoy"] = float(yoy)
        except Exception: pass
        tail = s.iloc[-min(len(s), 36):]
        try:
            x = np.arange(len(tail), dtype=float); b = np.polyfit(x, tail.values.astype(float), 1)[0]
            out["slope_ann"] = float(b * self._infer_months_in_year(s))
        except Exception: pass
        return out

    def _detect_turning_point(self, s: pd.Series) -> dict | None:
        s = s.dropna()
        if len(s) < 8: return None
        dif = s.diff().dropna()
        now = float(dif.iloc[-3:].mean())
        prev = float(dif.iloc[-6:-3].mean()) if len(dif) >= 6 else float(dif.iloc[:-3].mean())
        flip = (now > 0 and prev < 0) or (now < 0 and prev > 0)
        return {"flip": flip, "from": ("up" if prev > 0 else "down"), "to": ("up" if now > 0 else "down")} if flip else {"flip": False}

    def _zscore_flag(self, s: pd.Series, window: int = 36) -> dict | None:
        s = s.dropna()
        if len(s) < max(12, window // 2): return None
        tail = s.iloc[-min(len(s), window):]; mu, sd = float(tail.mean()), float(tail.std(ddof=0))
        if sd <= 0 or not np.isfinite(sd): return None
        z = (float(tail.iloc[-1]) - mu) / sd
        return {"z": float(z), "window": len(tail)}

    # -------------------- Export + Filters --------------------
    def _export_csv(self):
        if self.current_df is None or self.current_df.empty:
            messagebox.showinfo("Info", "Nothing to export yet."); return
        df = self.current_df.copy()
        path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV","*.csv")], title="Save data as CSV")
        if not path: return
        try:
            if not isinstance(df.index, pd.DatetimeIndex):
                df.index = pd.to_datetime(df.index, errors="coerce")
            df.to_csv(path, index_label="date"); self._set_status(f"Saved CSV → {path}")
        except Exception as e: self._error(e)

    def _export_ppt(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except Exception:
            messagebox.showinfo("Missing dependency",
                                "PowerPoint export needs 'python-pptx'.\nInstall:\n  pip install python-pptx"); return
        if not self.canvas:
            messagebox.showinfo("Info", "No chart to export yet."); return

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        self.canvas.figure.savefig(tmp.name, dpi=144, bbox_inches="tight")

        prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
        TF_RED = RGBColor(226, 0, 26); SLATE = RGBColor(45, 55, 72)

        title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        fill = title_box.fill; fill.solid(); fill.fore_color.rgb = TF_RED; title_box.line.fill.background()
        p = title_box.text_frame.paragraphs[0]; p.text = "Labor Market Dashboard"
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)

        slide.shapes.add_picture(tmp.name, Inches(0.5), Inches(1.2), height=Inches(4.5))

        insight_str = self._collect_insights_text()
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.3), Inches(2.3))
        tf = tx_box.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; p.text = "Insights"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.LEFT
        for line in insight_str.splitlines():
            if not line.strip(): continue
            q = tf.add_paragraph(); q.text = line; q.level = 1; q.font.size = Pt(12); q.font.color.rgb = SLATE

        path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                            filetypes=[("PowerPoint","*.pptx")],
                                            title="Save dashboard slide")
        if not path: return
        try: prs.save(path); self._set_status(f"Saved PowerPoint → {path}")
        except Exception as e: self._error(e)

    def _collect_insights_text(self) -> str:
        self.insights_text.configure(state="normal")
        txt = (self.insights_text.get("1.0","end") or "").strip() or "No insights available."
        self.insights_text.configure(state="disabled")
        return txt

    def _clear_filters(self):
        self.q.set(""); self.start_var.set(""); self.end_var.set(""); self.range_var.set("Max")
        self.normalize_var.set(False); self.yoy_var.set(False)
        self.forecast_var.set(False); self.forecast_h_var.set("12m"); self.scenario_var.set("Baseline")
        self.category_var.set(""); self.category_cb.set(""); self.subcategory_var.set(""); self.subcategory_cb.set("")
        self.dual_axis_var.set(False); self.dual_series_var.set("(none)"); self.dual_series_cb.configure(state="disabled")
        self._clear_results(); self._set_status("Filters cleared.")

    # -------------------- Treeview helpers --------------------
    def _wrap_title(self, text: str, width: int = 50) -> str:
        text = (text or "").strip()
        return "\n".join(textwrap.wrap(text, width=width)) if text else ""

    def _refresh_results_rowheight(self, style_name: str = "Results.Treeview", base_pad_px: int = 6, max_lines: int = 5):
        longest = 1
        for iid in self.results_tree.get_children(""):
            vals = self.results_tree.item(iid, "values")
            if not vals: continue
            title = str(vals[1]) if len(vals) > 1 else ""
            lines = title.count("\n") + 1 if title else 1
            if lines > longest: longest = lines
        longest = min(longest, max_lines)
        try:
            f = tkfont.nametofont("TkDefaultFont"); line_px = f.metrics("linespace")
            row_h = max(28, int(longest * (line_px + base_pad_px)))
        except Exception:
            row_h = 24 + (longest - 1) * 16
        ttk.Style(self).configure(style_name, rowheight=row_h)

    # -------------------- Helpers & Close --------------------
    def _bg(self, func, on_done):
        try: res = func(); self.after(0, lambda: on_done(res))
        except Exception as e: self.after(0, lambda: self._error(e))

    def _set_status(self, text: str):
        self.status_var.set(text); self.update_idletasks()

    def _error(self, e: Exception):
        self._set_status(f"Error: {e}"); messagebox.showerror("Error", str(e))

    def _on_close(self):
        try: plt.close("all")
        except Exception: pass
        try: self.master.quit(); self.master.destroy()
        except Exception: pass
        import sys; sys.exit(0)
